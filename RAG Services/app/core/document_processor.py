"""
Document Processing - Extract text from various formats
All FREE libraries
With Vietnamese optimization support
"""
import os
from pathlib import Path
from typing import List, Dict, Any
import re

# PDF processing
from pypdf import PdfReader

# Office documents
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook

# Markdown
import markdown

from .config import settings
from .vietnamese_processor import get_vietnamese_processor

class DocumentProcessor:
    """
    Process various document formats - all FREE tools
    Supported: PDF, DOCX, PPTX, XLSX, TXT, MD, HTML
    """
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF using pypdf (FREE)"""
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n\n"
            return text.strip()
        except Exception as e:
            print(f"❌ Error reading PDF: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from DOCX using python-docx (FREE)"""
        try:
            doc = DocxDocument(file_path)
            text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            return text
        except Exception as e:
            print(f"❌ Error reading DOCX: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PPTX using python-pptx (FREE)"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text.strip()
        except Exception as e:
            print(f"❌ Error reading PPTX: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_xlsx(file_path: str) -> str:
        """Extract text from XLSX using openpyxl (FREE)"""
        try:
            wb = load_workbook(file_path, data_only=True)
            text = ""
            for sheet in wb.worksheets:
                text += f"Sheet: {sheet.title}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
            return text.strip()
        except Exception as e:
            print(f"❌ Error reading XLSX: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """Extract text from TXT"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception as e:
                print(f"❌ Error reading TXT: {e}")
                return ""
    
    @staticmethod
    def extract_text_from_md(file_path: str) -> str:
        """Extract text from Markdown"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_text = f.read()
            # Convert to HTML then strip tags for plain text
            html = markdown.markdown(md_text)
            text = re.sub('<[^<]+?>', '', html)
            return text
        except Exception as e:
            print(f"❌ Error reading MD: {e}")
            return ""
    
    @classmethod
    def process_file(cls, file_path: str) -> str:
        """
        Process any supported file format
        
        Args:
            file_path: Path to file
            
        Returns:
            Extracted text
        """
        ext = Path(file_path).suffix.lower()
        
        extractors = {
            '.pdf': cls.extract_text_from_pdf,
            '.docx': cls.extract_text_from_docx,
            '.doc': cls.extract_text_from_docx,
            '.pptx': cls.extract_text_from_pptx,
            '.xlsx': cls.extract_text_from_xlsx,
            '.txt': cls.extract_text_from_txt,
            '.md': cls.extract_text_from_md,
            '.html': cls.extract_text_from_txt,
        }
        
        extractor = extractors.get(ext)
        if extractor:
            print(f"📄 Processing {ext} file: {Path(file_path).name}")
            return extractor(file_path)
        else:
            print(f"⚠️  Unsupported file type: {ext}")
            return ""
    
    @staticmethod
    def chunk_text(
        text: str, 
        chunk_size: int = None, 
        overlap: int = None,
        use_vietnamese: bool = True
    ) -> List[str]:
        """
        Split text into chunks with overlap
        Vietnamese-aware sentence segmentation
        
        Args:
            text: Text to chunk
            chunk_size: Size of each chunk (default from settings)
            overlap: Overlap between chunks (default from settings)
            use_vietnamese: Use Vietnamese-aware chunking
            
        Returns:
            List of text chunks
        """
        chunk_size = chunk_size or settings.CHUNK_SIZE
        overlap = overlap or settings.CHUNK_OVERLAP
        
        # Try Vietnamese-aware chunking first
        if use_vietnamese:
            try:
                vi_processor = get_vietnamese_processor()
                
                # Detect if text is Vietnamese
                lang = vi_processor.detect_language(text)
                
                if lang == 'vi':
                    print(f"   🇻🇳 Vietnamese text detected - using optimized chunking")
                    return vi_processor.chunk_vietnamese_text(text, chunk_size, overlap)
            except Exception as e:
                print(f"   ⚠️  Vietnamese chunking failed, using fallback: {e}")
        
        # Fallback to simple sentence-based chunking
        sentences = re.split(r'(?<=[.!?])\s+', text)
        
        chunks = []
        current_chunk = []
        current_size = 0
        
        for sentence in sentences:
            sentence_size = len(sentence.split())
            
            if current_size + sentence_size > chunk_size and current_chunk:
                # Save current chunk
                chunks.append(' '.join(current_chunk))
                
                # Start new chunk with overlap
                overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                current_chunk = overlap_text + [sentence]
                current_size = sum(len(s.split()) for s in current_chunk)
            else:
                current_chunk.append(sentence)
                current_size += sentence_size
        
        # Add last chunk
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks
    
    @classmethod
    def process_and_chunk(
        cls,
        file_path: str,
        chunk_size: int = None,
        overlap: int = None,
        use_vietnamese: bool = True
    ) -> List[Dict[str, Any]]:
        """
        Process file and return chunks with metadata
        With Vietnamese text preprocessing
        
        Args:
            file_path: Path to file
            chunk_size: Chunk size
            overlap: Overlap size
            use_vietnamese: Use Vietnamese optimization
            
        Returns:
            List of dicts with 'text' and 'metadata'
        """
        # Extract text
        text = cls.process_file(file_path)
        if not text:
            return []
        
        # Preprocess text if Vietnamese
        if use_vietnamese:
            try:
                vi_processor = get_vietnamese_processor(
                    use_tokenization=False,  # Don't tokenize at this stage
                    remove_stopwords=False
                )
                
                # Clean and normalize
                text = vi_processor.clean_text(text)
                
                # Get statistics
                stats = vi_processor.get_statistics(text)
                if stats['vietnamese_ratio'] > 0.1:
                    print(f"   🇻🇳 Vietnamese content: {stats['vietnamese_ratio']*100:.1f}%")
            except Exception as e:
                print(f"   ⚠️  Vietnamese preprocessing failed: {e}")
        
        # Chunk text
        chunks = cls.chunk_text(text, chunk_size, overlap, use_vietnamese)
        
        # Add metadata
        file_name = Path(file_path).name
        result = []
        for i, chunk in enumerate(chunks):
            result.append({
                'text': chunk,
                'metadata': {
                    'source': file_name,
                    'chunk_id': i,
                    'total_chunks': len(chunks),
                    'file_type': Path(file_path).suffix.lower()
                }
            })
        
        print(f"✅ Created {len(chunks)} chunks from {file_name}")
        return result
